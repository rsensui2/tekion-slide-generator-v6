#!/usr/bin/env python3
"""TEKION Slide Generator v6 - 既存デッキ取り込み（Phase 8 拡張）

画像ベースの PPTX（各スライドが1枚の画像で構成されたデッキ）や単体画像を
セッションに取り込み、manifest に登録する。取り込んだスライドは通常の
生成スライドと同様に、スライドダッシュボードでのレビュー・差分編集・export の対象になる。

「既存のパワポをこのスキルで直していく」ための入り口:
    1. python3 import_deck.py --session-dir <新規/既存セッション> --file deck.pptx
    2. python3 review_deck.py --session-dir <セッション> --serve
    3. 赤入れ → 差分編集 → PPTX/PDF 書き出し

対応形式:
  - .pptx : 各スライドから最大面積の画像を抽出。画像が無いスライド（図形・テキスト
            ボックスで組まれたネイティブスライド）は、LibreOffice(soffice) があれば
            PDF 経由でページ全体をレンダリングして取り込む
  - .pdf  : 各ページを画像としてレンダリング（pymupdf）
  - .png / .jpg / .jpeg / .webp : 1ファイル = 1スライドとして取り込み
"""
from __future__ import annotations

import argparse
import os
import re
import sys
from io import BytesIO

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from manifest_utils import load_manifest, save_manifest, update_entry, validate_image

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}


def _sanitize_stem(name: str) -> str:
    stem = os.path.splitext(os.path.basename(name))[0]
    stem = re.sub(r"[\\/:*?\"<>|\s]+", "_", stem).strip("_")
    return stem or "imported"


def _unique_base(slides: dict, base: str) -> str:
    if base not in slides:
        return base
    n = 2
    while f"{base}-{n}" in slides:
        n += 1
    return f"{base}-{n}"


def _save_as_png(blob: bytes, out_path: str) -> None:
    """任意形式の画像バイト列を PNG として保存する（RGBA→RGB 白合成）。"""
    from PIL import Image
    img = Image.open(BytesIO(blob))
    img.load()
    if img.mode in ("RGBA", "LA", "P"):
        img = img.convert("RGBA")
        bg = Image.new("RGB", img.size, (255, 255, 255))
        bg.paste(img, mask=img.split()[-1])
        img = bg
    elif img.mode != "RGB":
        img = img.convert("RGB")
    img.save(out_path, format="PNG")


def _register(manifest: dict, slide_base: str, image_path: str) -> None:
    update_entry(manifest, slide_base,
                 state="validated",
                 current_image=os.path.abspath(image_path),
                 raw_image=None,  # 取り込み画像に raw は無い（編集時は焼き込み済みから編集）
                 versions=[os.path.abspath(image_path)],
                 prompt_sha256=None,
                 imported=True)


def import_pptx(pptx_path: str, session_dir: str) -> tuple[list[str], list[int]]:
    """PPTX の各スライドから最大面積の画像を抽出して取り込む。

    Returns: (取り込んだ slide_base のリスト, 画像が見つからずスキップしたスライド番号)
    """
    from pptx import Presentation

    session_dir = os.path.abspath(session_dir)
    images_dir = os.path.join(session_dir, "images")
    os.makedirs(images_dir, exist_ok=True)
    manifest_path = os.path.join(session_dir, "manifest.json")
    manifest = load_manifest(manifest_path)

    prs = Presentation(pptx_path)
    stem = _sanitize_stem(pptx_path)
    added, skipped = [], []

    for idx, slide in enumerate(prs.slides, start=1):
        best_blob, best_area = None, -1
        for shape in slide.shapes:
            # PICTURE(13) / LINKED_PICTURE(14) を対象。画像を持つ shape のみ
            image = getattr(shape, "image", None)
            if image is None:
                continue
            area = int(shape.width or 0) * int(shape.height or 0)
            if area > best_area:
                best_area = area
                best_blob = image.blob
        if best_blob is None:
            skipped.append(idx)
            continue

        base = _unique_base(manifest.get("slides", {}), f"{stem}_{idx:02d}")
        out_path = os.path.join(images_dir, f"{base}.png")
        _save_as_png(best_blob, out_path)
        problem = validate_image(out_path)
        if problem:
            print(f"⚠️  slide {idx}: 検証失敗のためスキップ ({problem})", file=sys.stderr)
            os.unlink(out_path)
            skipped.append(idx)
            continue
        _register(manifest, base, out_path)
        added.append(base)

    save_manifest(manifest_path, manifest)
    return added, skipped


def import_pdf(pdf_path: str, session_dir: str) -> tuple[list[str], list[int]]:
    """PDF の各ページを画像にレンダリングして取り込む（pymupdf）。"""
    try:
        import fitz  # pymupdf
    except ImportError:
        raise ValueError("PDF の取り込みには pymupdf が必要です: pip install pymupdf")

    session_dir = os.path.abspath(session_dir)
    images_dir = os.path.join(session_dir, "images")
    os.makedirs(images_dir, exist_ok=True)
    manifest_path = os.path.join(session_dir, "manifest.json")
    manifest = load_manifest(manifest_path)

    stem = _sanitize_stem(pdf_path)
    added, skipped = [], []
    with fitz.open(pdf_path) as doc:
        for idx, page in enumerate(doc, start=1):
            # 長辺 ~2560px 相当でレンダリング（2K スライドと同等の精細度）
            rect = page.rect
            zoom = 2560 / max(rect.width, rect.height)
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
            base = _unique_base(manifest.get("slides", {}), f"{stem}_{idx:02d}")
            out_path = os.path.join(images_dir, f"{base}.png")
            pix.save(out_path)
            problem = validate_image(out_path)
            if problem:
                print(f"⚠️  page {idx}: 検証失敗のためスキップ ({problem})", file=sys.stderr)
                os.unlink(out_path)
                skipped.append(idx)
                continue
            _register(manifest, base, out_path)
            added.append(base)

    save_manifest(manifest_path, manifest)
    return added, skipped


def _find_soffice() -> str | None:
    import shutil as _shutil
    found = _shutil.which("soffice")
    if found:
        return found
    candidates = ["/Applications/LibreOffice.app/Contents/MacOS/soffice"]
    for env_key in ("ProgramFiles", "ProgramFiles(x86)"):
        base = os.environ.get(env_key)
        if base:
            candidates.append(os.path.join(base, "LibreOffice", "program", "soffice.exe"))
    for path in candidates:
        if os.path.exists(path):
            return path
    return None


def _pptx_to_pdf(pptx_path: str, out_dir: str) -> str | None:
    """LibreOffice で PPTX を PDF に変換する（ネイティブスライドの図化）。"""
    import subprocess
    soffice = _find_soffice()
    if not soffice:
        return None
    os.makedirs(out_dir, exist_ok=True)
    result = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
        capture_output=True, encoding="utf-8", errors="replace", timeout=180)
    pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    return pdf_path if result.returncode == 0 and os.path.exists(pdf_path) else None


def import_image(image_path: str, session_dir: str) -> str:
    """単体画像を1スライドとして取り込む。Returns: slide_base"""
    session_dir = os.path.abspath(session_dir)
    images_dir = os.path.join(session_dir, "images")
    os.makedirs(images_dir, exist_ok=True)
    manifest_path = os.path.join(session_dir, "manifest.json")
    manifest = load_manifest(manifest_path)

    stem = _sanitize_stem(image_path)
    base = _unique_base(manifest.get("slides", {}), f"{stem}_01")
    out_path = os.path.join(images_dir, f"{base}.png")
    with open(image_path, "rb") as f:
        _save_as_png(f.read(), out_path)
    problem = validate_image(out_path)
    if problem:
        os.unlink(out_path)
        raise ValueError(f"画像の検証に失敗: {problem}")
    _register(manifest, base, out_path)
    save_manifest(manifest_path, manifest)
    return base


def import_file(path: str, session_dir: str) -> dict:
    """拡張子で振り分けて取り込む。Returns: {added: [...], skipped: [...], note: str}"""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        added, skipped = import_pptx(path, session_dir)
        # 画像が抽出できないスライドがある = 図形・テキストで組まれたネイティブデッキ。
        # LibreOffice があれば PDF 経由でページ全体を図化して取り込み直す
        if skipped and not added:
            pdf = _pptx_to_pdf(path, os.path.join(session_dir, "imports"))
            if pdf:
                added, skipped = import_pdf(pdf, session_dir)
                return {"added": added, "skipped": skipped,
                        "note": "画像抽出できないデッキのため LibreOffice で図化して取り込み"}
            return {"added": added, "skipped": skipped,
                    "note": "画像が抽出できません。PowerPoint から PDF で書き出して取り込んでください"}
        return {"added": added, "skipped": skipped, "note": ""}
    if ext == ".pdf":
        added, skipped = import_pdf(path, session_dir)
        return {"added": added, "skipped": skipped, "note": ""}
    if ext in IMAGE_EXTS:
        return {"added": [import_image(path, session_dir)], "skipped": [], "note": ""}
    raise ValueError(f"未対応の形式: {ext}（対応: .pptx / .pdf / {', '.join(sorted(IMAGE_EXTS))}）")


def main() -> int:
    ap = argparse.ArgumentParser(description="Import an image-based PPTX or images into a session (v6)")
    ap.add_argument("--session-dir", required=True,
                    help="取り込み先セッション（無ければ作成される）")
    ap.add_argument("--file", required=True, action="append",
                    help="取り込むファイル（.pptx / 画像。複数指定可）")
    args = ap.parse_args()

    os.makedirs(args.session_dir, exist_ok=True)
    total_added, total_skipped = [], 0
    for path in args.file:
        if not os.path.exists(path):
            print(f"❌ ファイルが見つかりません: {path}")
            return 1
        result = import_file(path, args.session_dir)
        total_added.extend(result["added"])
        total_skipped += len(result["skipped"])
        print(f"✅ {os.path.basename(path)}: {len(result['added'])}枚を取り込み"
              + (f"（{len(result['skipped'])}枚は画像なしでスキップ）" if result["skipped"] else ""))

    print(f"\n合計 {len(total_added)}枚を取り込みました → {args.session_dir}/manifest.json")
    if total_added:
        print("次: スライドダッシュボード（ハブ or review_deck.py --serve）で赤入れ → 差分編集で直せます")
    if total_skipped:
        print(f"⚠️  {total_skipped}枚は画像が抽出できずスキップ（図形ベースのスライドは対象外）")
    return 0


if __name__ == "__main__":
    sys.exit(main())
